#!/usr/bin/env python3
"""
Génère une présentation PowerPoint à partir du rapport
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def add_title_slide(prs, title, subtitle):
    """Ajoute une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, content_lines):
    """Ajoute une slide avec titre et contenu en bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for line in content_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)

def add_image_slide(prs, title, image_path, caption=""):
    """Ajoute une slide avec titre et image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Titre
    title_shape = slide.shapes.title
    title_shape.text = title

    # Image centrée
    if Path(f"images/{image_path}").exists():
        left = Inches(2)
        top = Inches(2)
        height = Inches(4)
        slide.shapes.add_picture(f"images/{image_path}", left, top, height=height)

        if caption:
            # Légende sous l'image
            text_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.text = caption
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_presentation():
    """Crée la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Slide de titre
    add_title_slide(
        prs,
        "Morpion Q-Learning",
        "Master 2 Data Science - Ynov\nModule: Machine Learning\nJanvier 2026"
    )

    # 2. Contexte
    add_content_slide(
        prs,
        "1. Contexte",
        [
            "Jeu: Morpion (Tic-Tac-Toe) sur grille 3×3",
            "Objectif: Créer un agent intelligent via Q-Learning",
            "Cible: >85% de victoires contre adversaire aléatoire",
            "Méthode: Apprentissage par renforcement sans programmation explicite"
        ]
    )

    # 3. Q-Learning
    add_content_slide(
        prs,
        "2. Q-Learning - Principe",
        [
            "Équation: Q(s,a) ← Q(s,a) + α[r + γ max Q(s',a') - Q(s,a)]",
            "Hyperparamètres:",
            "  • α = 0.2 (taux d'apprentissage)",
            "  • γ = 0.95 (discount factor)",
            "  • ε = 1.0 → 0.05 (exploration/exploitation)",
            "",
            "Stratégie ε-greedy:",
            "  • Avec probabilité ε: action aléatoire (exploration)",
            "  • Avec probabilité 1-ε: meilleure action (exploitation)"
        ]
    )

    # 4. Architecture
    add_content_slide(
        prs,
        "3. Architecture du projet",
        [
            "• main.py - Point d'entrée de l'application",
            "• gui.py - Interface Pygame moderne",
            "• environment.py - Règles du jeu et récompenses",
            "• agent.py - Algorithme Q-Learning",
            "• trainer.py - Entraînement symétrique X/O",
            "• agent_trained.pkl - Modèle sauvegardé",
            "",
            "État: Tuple de 9 valeurs (0=vide, 1=X, 2=O)",
            "Récompenses: +1 victoire, -1 défaite, 0 nul"
        ]
    )

    # 5. Interface - Image menu
    add_image_slide(
        prs,
        "4. Interface graphique",
        "image1.png",
        "Menu principal avec 3 modes de jeu + Entraînement/Évaluation"
    )

    # 6. Interface - Image partie
    add_image_slide(
        prs,
        "4. Interface graphique (suite)",
        "image6.png",
        "Partie en cours - X (orange) vs O (vert)"
    )

    # 7. Implémentations réalisées
    add_content_slide(
        prs,
        "5. Implémentations réalisées",
        [
            "✅ Environnement de jeu complet",
            "✅ Agent Q-Learning avec stratégie ε-greedy",
            "✅ Entraînement symétrique (50% X, 50% O)",
            "✅ Interface graphique moderne (Pygame)",
            "✅ 3 modes: Humain vs Humain, Humain vs IA, IA vs IA",
            "✅ Sauvegarde/chargement de l'agent",
            "✅ Module d'évaluation quantitative"
        ]
    )

    # 8. Difficultés rencontrées
    add_content_slide(
        prs,
        "6. Difficultés rencontrées",
        [
            "Problème 1: Manque de temps d'entraînement",
            "  → Agent sous-performant (absence hier)",
            "",
            "Problème 2: Asymétrie X/O",
            "  → Solution: Entraînement symétrique 50/50",
            "",
            "Problème 3: Convergence lente (10k épisodes → 70%)",
            "  → Solution: 100k épisodes + ajustement hyperparamètres",
            "",
            "Problème 4: Incohérences visuelles",
            "  → Solution: Anti-aliasing + correction design"
        ]
    )

    # 9. Résultats théoriques
    add_content_slide(
        prs,
        "7. Résultats théoriques",
        [
            "Évolution avec entraînement complet:",
            "",
            "  • 0 épisodes → ~50% (aléatoire)",
            "  • 10 000 épisodes → ~75% (stratégies basiques)",
            "  • 30 000 épisodes → 85-90% (stratégie mature)",
            "  • 100 000 épisodes → 89-92% (quasi-optimal)",
            "",
            "Limitation actuelle:",
            "Agent non optimal par manque de temps d'entraînement"
        ]
    )

    # 10. Image évaluation
    add_image_slide(
        prs,
        "7. Résultats - Évaluation",
        "image5.png",
        "Résultats d'évaluation après entraînement complet: 89.5% victoires"
    )

    # 11. Axes d'amélioration
    add_content_slide(
        prs,
        "8. Axes d'amélioration",
        [
            "1. Plus d'entraînement",
            "   → Laisser tourner plusieurs heures (85-90%)",
            "",
            "2. Deep Q-Learning",
            "   → Réseau de neurones au lieu de Q-table",
            "",
            "3. Adversaire intelligent",
            "   → Entraîner contre Minimax",
            "",
            "4. Grilles plus grandes (4×4, 5×5)",
            "",
            "5. Interface enrichie (animations, sons)"
        ]
    )

    # 12. Conclusion
    add_content_slide(
        prs,
        "9. Conclusion",
        [
            "✓ Q-Learning fonctionne pour le morpion",
            "✓ Entraînement symétrique crucial (évite biais X/O)",
            "✓ Interface complète et fonctionnelle",
            "",
            "Concepts clés appris:",
            "  • Équilibre exploration/exploitation (ε-greedy)",
            "  • Importance du temps d'entraînement",
            "  • Hyperparamètres critiques pour convergence",
            "",
            "Avec entraînement complet: 89-92% victoires possibles"
        ]
    )

    # Sauvegarder
    output_file = "presentation_morpion.pptx"
    prs.save(output_file)
    print(f"[OK] Presentation creee: {output_file}")
    print(f"[OK] {len(prs.slides)} slides generees")

if __name__ == "__main__":
    create_presentation()
